import fitz  # PyMuPDF
import docx
from pptx import Presentation
from PIL import Image
import io

def extract_from_pdf(filepath):
    """
    Extracts images of each page from a PDF.
    Yields: (PIL.Image, True)
    """
    doc = fitz.open(filepath)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        # 150 DPI is usually good enough for OCR
        pix = page.get_pixmap(dpi=150)
        
        # Convert to PIL Image
        img_bytes = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_bytes))
        
        yield (img, True)

def extract_from_pptx(filepath):
    """
    Extracts text from each slide in a PPTX.
    Yields: (str, False)
    """
    prs = Presentation(filepath)
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        full_text = "\n".join(slide_text).strip()
        if full_text:
             yield (full_text, False)
        # Even if empty, we could yield to maintain slide numbers, but skipping empty slides is usually better.

def extract_from_docx(filepath, words_per_chunk=500):
    """
    Extracts text from a Word document and chunks it since Word doesn't have native pagination.
    Yields: (str, False)
    """
    doc = docx.Document(filepath)
    
    current_chunk = []
    current_word_count = 0
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        words = text.split()
        current_chunk.append(text)
        current_word_count += len(words)
        
        if current_word_count >= words_per_chunk:
            yield ("\n\n".join(current_chunk), False)
            current_chunk = []
            current_word_count = 0
            
    if current_chunk:
        yield ("\n\n".join(current_chunk), False)

def extract_from_markdown(filepath):
    """
    Extracts text from a Markdown file.
    If the file already contains '---' as dividers, it splits them directly.
    Otherwise, it imports and calls `divide_markdown_with_llm` from ollama_client to 
    insert the dividers smartly, then splits it.
    Yields: (str, False)
    """
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # If the file already has horizontal rules (standard or alternate syntaxes)
    if "\n---\n" in content or "\n***\n" in content or "\n___\n" in content:
        # Standardize split logic for simplicity
        normalized_content = content.replace("\n***\n", "\n---\n").replace("\n___\n", "\n---\n")
        chunks = normalized_content.split("\n---\n")
    else:
        print("  No separators found in Markdown. Using LLM to intelligently divide the document into sections...")
        # Local import to avoid circular dependency since ollama_client might import from extractor later or vice versa
        from ollama_client import divide_markdown_with_llm
        divided_content = divide_markdown_with_llm(content)
        # LLM might use variations, so normalize before split
        normalized_content = divided_content.replace("\n***\n", "\n---\n").replace("\n___\n", "\n---\n")
        chunks = normalized_content.split("\n---\n")
        
    for chunk in chunks:
        chunk = chunk.strip()
        if chunk:
            yield (chunk, False)
